import streamlit as st
import os
import time
import uuid
from datetime import datetime, timedelta
from azure.storage.blob import BlobServiceClient, generate_blob_sas, BlobSasPermissions, generate_container_sas, ContainerSasPermissions
from azure.ai.translation.document import DocumentTranslationClient, DocumentTranslationInput, TranslationTarget
from azure.core.credentials import AzureKeyCredential
import urllib.parse

import requests
import json
from pathlib import Path

# -----------------------------
# 설정 및 비밀 관리
# -----------------------------
st.set_page_config(page_title="Azure 문서 번역기", page_icon="🌏", layout="centered")

def get_secret(key):
    if key in st.secrets:
        return st.secrets[key]
    return os.environ.get(key)

# 필수 자격 증명
STORAGE_CONN_STR = get_secret("AZURE_STORAGE_CONNECTION_STRING")
TRANSLATOR_KEY = get_secret("AZURE_TRANSLATOR_KEY")
TRANSLATOR_ENDPOINT = get_secret("AZURE_TRANSLATOR_ENDPOINT")
CONTAINER_NAME = get_secret("AZURE_BLOB_CONTAINER_NAME") or "blob-leesunguk"

# -----------------------------
# Azure 클라이언트 헬퍼
# -----------------------------
def get_blob_service_client():
    if not STORAGE_CONN_STR:
        st.error("Azure Storage Connection String이 설정되지 않았습니다.")
        st.stop()
    return BlobServiceClient.from_connection_string(STORAGE_CONN_STR)

def get_translation_client():
    if not TRANSLATOR_KEY or not TRANSLATOR_ENDPOINT:
        st.error("Azure Translator Key 또는 Endpoint가 설정되지 않았습니다.")
        st.stop()
    return DocumentTranslationClient(TRANSLATOR_ENDPOINT, AzureKeyCredential(TRANSLATOR_KEY))

def generate_sas_url(blob_service_client, container_name, blob_name=None, permission="r", expiry_hours=1):
    """
    Blob 또는 Container에 대한 SAS URL 생성
    blob_name이 있으면 Blob SAS, 없으면 Container SAS (Write용)
    """
    import urllib.parse
    
    account_name = blob_service_client.account_name
    
    # Connection String으로 생성된 경우 credential은 dict일 수 있음
    if hasattr(blob_service_client.credential, 'account_key'):
        account_key = blob_service_client.credential.account_key
    else:
        account_key = blob_service_client.credential['account_key']
    
    # 시계 오차(Clock Skew) 방지를 위해 시작 시간을 15분 전으로 설정
    start = datetime.utcnow() - timedelta(minutes=15)
    expiry = datetime.utcnow() + timedelta(hours=expiry_hours)
    
    # 항상 Container SAS를 사용 (Source/Target 모두 더 안정적)
    # Source의 경우 Read/List, Target의 경우 Write/List/Read 필요
    # 편의상 모든 권한을 부여한 Container SAS 하나로 통일하거나, 구분 가능
    # 여기서는 구분 없이 Container 수준의 강력한 SAS를 발급하여 오류 가능성 차단
    
    sas_token = generate_container_sas(
        account_name=account_name,
        container_name=container_name,
        account_key=account_key,
        permission=ContainerSasPermissions(write=True, list=True, read=True, delete=True),
        start=start,
        expiry=expiry
    )
    
    base_url = f"https://{account_name}.blob.core.windows.net/{container_name}"
    
    if blob_name:
        # Blob 경로가 있는 경우 URL에 추가 (SAS는 컨테이너 레벨이라 서명 불일치 없음)
        encoded_blob_name = urllib.parse.quote(blob_name, safe='/')
        return f"{base_url}/{encoded_blob_name}?{sas_token}"
    else:
        # 컨테이너 루트 URL
        return f"{base_url}?{sas_token}"
# -----------------------------
# 작업 상태 관리 (Local JSON)
# -----------------------------
JOBS_FILE = "jobs.json"

def load_jobs():
    if os.path.exists(JOBS_FILE):
        try:
            with open(JOBS_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except:
            return {}
    return {}

def save_job(job_id, job_data):
    jobs = load_jobs()
    jobs[job_id] = job_data
    with open(JOBS_FILE, "w", encoding="utf-8") as f:
        json.dump(jobs, f, ensure_ascii=False, indent=2)

def update_job_status(job_id, status):
    jobs = load_jobs()
    if job_id in jobs:
        jobs[job_id]["status"] = status
        with open(JOBS_FILE, "w", encoding="utf-8") as f:
            json.dump(jobs, f, ensure_ascii=False, indent=2)

# -----------------------------
# -----------------------------
# UI 구성
# -----------------------------
st.title("🌏 Azure 문서 번역기")
st.caption("Azure Document Translation & Blob Storage 기반")

# 지원 언어 목록 가져오기 (API)
@st.cache_data
def get_supported_languages():
    try:
        url = "https://api.cognitive.microsofttranslator.com/languages?api-version=3.0&scope=translation"
        # Accept-Language 헤더를 'ko'로 설정하여 언어 이름을 한국어로 받음
        headers = {"Accept-Language": "ko"}
        response = requests.get(url, headers=headers)
        response.raise_for_status()
        data = response.json()
        
        languages = {}
        for code, info in data['translation'].items():
            # "한국어 이름 (원어 이름)" 형식으로 표시 (예: 영어 (English))
            label = f"{info['name']} ({info['nativeName']})"
            languages[label] = code
        return languages
    except Exception as e:
        st.error(f"언어 목록을 가져오는데 실패했습니다: {e}")
        # 실패 시 기본 언어 제공
        return {"한국어 (Korean)": "ko", "영어 (English)": "en"}

LANGUAGES = get_supported_languages()

# 언어 코드별 파일 접미사 매핑 (기본적으로 대문자 코드를 사용하되, 일부 커스텀 가능)
# 여기서는 자동 생성 로직을 사용하므로 별도 딕셔너리 불필요, 
# 다만 중국어 등 특수 케이스를 위해 남겨둘 수 있음.
LANG_SUFFIX_OVERRIDE = {
    "zh-Hans": "CN",
    "zh-Hant": "TW",
}

with st.sidebar:
    st.header("메뉴")
    menu = st.radio("이동", ["번역하기", "작업 상태", "파일 보관함"])
    
    st.divider()
    
    if menu == "번역하기":
        st.header("설정")
        # 한국어를 기본값으로 찾기
        default_index = 0
        lang_labels = list(LANGUAGES.keys())
        for i, label in enumerate(lang_labels):
            if "Korean" in label or "한국어" in label:
                default_index = i
                break
                
        target_lang_label = st.selectbox("목표 언어 선택", lang_labels, index=default_index)
        target_lang_code = LANGUAGES[target_lang_label]
        st.info(f"선택된 목표 언어: {target_lang_code}")

    # 자격 증명 상태 확인
    if STORAGE_CONN_STR and TRANSLATOR_KEY:
        st.success("✅ Azure 자격 증명 확인됨")
    else:
        st.warning("⚠️ Azure 자격 증명이 누락되었습니다. secrets.toml을 확인하세요.")

if menu == "번역하기":
    uploaded_file = st.file_uploader("번역할 문서 업로드 (PPTX, PDF, DOCX, XLSX 등)", type=["pptx", "pdf", "docx", "xlsx"])

    if st.button("번역 시작", type="primary", disabled=not uploaded_file):
        if not uploaded_file:
            st.error("파일을 업로드해주세요.")
        else:
            with st.spinner("Azure Blob에 파일 업로드 중..."):
                try:
                    blob_service_client = get_blob_service_client()
                    container_client = blob_service_client.get_container_client(CONTAINER_NAME)
                    
                    # 컨테이너 접근 권한 확인
                    try:
                        if not container_client.exists():
                            container_client.create_container()
                    except Exception as e:
                        if "AuthenticationFailed" in str(e):
                            st.error("🚨 인증 실패: Azure Storage Key가 올바르지 않습니다. Secrets 설정을 확인해주세요.")
                            st.stop()
                        else:
                            raise e

                    # 파일명 유니크하게 처리
                    file_uuid = str(uuid.uuid4())[:8]
                    original_filename = uploaded_file.name
                    input_blob_name = f"input/{file_uuid}/{original_filename}"
                    
                    # 업로드
                    blob_client = container_client.get_blob_client(input_blob_name)
                    blob_client.upload_blob(uploaded_file, overwrite=True)
                    
                    st.success("업로드 완료! 번역 요청 중...")
                    
                    # SAS 생성
                    source_url = generate_sas_url(blob_service_client, CONTAINER_NAME, input_blob_name)
                    
                    # Target URL 설정
                    target_base_url = f"https://{blob_service_client.account_name}.blob.core.windows.net/{CONTAINER_NAME}"
                    # Target URL은 컨테이너 또는 폴더 경로여야 함 (파일 경로 불가)
                    # 파일명 보존을 위해 폴더 경로 끝에 '/'를 반드시 붙여야 함
                    target_output_url = f"{target_base_url}/output/{file_uuid}/?{generate_sas_url(blob_service_client, CONTAINER_NAME).split('?')[1]}"
                    
                except Exception as e:
                    st.error(f"업로드/SAS 생성 실패: {e}")
                    st.stop()

            with st.spinner("번역 작업 요청 및 대기 중..."):
                try:
                    client = get_translation_client()
                    

                    poller = client.begin_translation(
                        inputs=[
                            DocumentTranslationInput(
                                source_url=source_url,
                                storage_type="File",
                                targets=[
                                    TranslationTarget(
                                        target_url=target_output_url,
                                        language=target_lang_code
                                    )
                                ]
                            )
                        ]
                    )
                    
                    job_id = poller.id
                    
                    # Save Job Info locally
                    job_data = {
                        "filename": original_filename,
                        "file_uuid": file_uuid,
                        "target_lang_code": target_lang_code,
                        "created_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                        "status": "Running" 
                    }
                    save_job(job_id, job_data)
                    
                    st.success(f"번역 작업이 시작되었습니다! (Job ID: {job_id})")
                    st.info("좌측 메뉴의 '작업 상태'에서 진행 상황을 확인하고 결과를 다운로드하세요.")
                    
                except Exception as e:
                    st.error(f"번역 요청 중 오류 발생: {e}")

elif menu == "작업 상태":
    st.subheader("⏳ 번역 작업 상태")
    
    jobs = load_jobs()
    if not jobs:
        st.info("기록된 작업이 없습니다.")
    else:
        # Sort by date desc
        sorted_job_ids = sorted(jobs.keys(), key=lambda x: jobs[x].get('created_at', ''), reverse=True)
        
        client = get_translation_client()
        blob_service_client = get_blob_service_client()
        container_client = blob_service_client.get_container_client(CONTAINER_NAME)

        for job_id in sorted_job_ids:
            job = jobs[job_id]
            
            with st.expander(f"{job.get('created_at')} - {job.get('filename')} ({job.get('status')})", expanded=True):
                col1, col2 = st.columns([3, 1])
                
                with col1:
                    st.write(f"**Job ID:** `{job_id}`")
                    st.write(f"**Target Lang:** {job.get('target_lang_code')}")
                
                with col2:
                    # Check Status Button
                    if st.button("상태 확인", key=f"check_{job_id}"):
                        try:
                            status_obj = client.get_translation_status(job_id)
                            new_status = status_obj.status
                            update_job_status(job_id, new_status)
                            st.rerun()
                        except Exception as e:
                            st.error(f"상태 확인 실패: {e}")

                # If Succeeded, Show Finalize/Download
                if job.get("status") == "Succeeded":
                    if st.button("결과 처리 및 다운로드", key=f"finalize_{job_id}", type="primary"):
                        with st.spinner("결과 파일 처리 중..."):
                            try:
                                file_uuid = job.get("file_uuid")
                                target_lang_code = job.get("target_lang_code")
                                
                                # 결과 파일 찾기
                                output_prefix_search = f"output/{file_uuid}"
                                output_blobs = list(container_client.list_blobs(name_starts_with=output_prefix_search))
                                
                                if not output_blobs:
                                    st.error("결과 파일을 찾을 수 없습니다.")
                                else:
                                    download_links = []
                                    for blob in output_blobs:
                                        blob_name = blob.name
                                        file_name = blob_name.split("/")[-1]
                                        
                                        # Rename Logic
                                        suffix = LANG_SUFFIX_OVERRIDE.get(target_lang_code, target_lang_code.upper())
                                        name_part, ext_part = os.path.splitext(file_name)
                                        
                                        final_blob_name = blob_name
                                        final_file_name = file_name
                                        
                                        if not name_part.endswith(f"_{suffix}"):
                                            new_file_name = f"{name_part}_{suffix}{ext_part}"
                                            new_blob_name = f"output/{file_uuid}/{new_file_name}"
                                            
                                            try:
                                                source_blob = container_client.get_blob_client(blob_name)
                                                dest_blob = container_client.get_blob_client(new_blob_name)
                                                source_sas = generate_sas_url(blob_service_client, CONTAINER_NAME, blob_name)
                                                dest_blob.start_copy_from_url(source_sas)
                                                
                                                for _ in range(10):
                                                    props = dest_blob.get_blob_properties()
                                                    if props.copy.status == "success":
                                                        break
                                                    time.sleep(0.2)
                                                
                                                source_blob.delete_blob()
                                                final_blob_name = new_blob_name
                                                final_file_name = new_file_name
                                            except Exception as e:
                                                st.warning(f"Rename failed: {e}")

                                        # PPTX Font Fix
                                        if final_file_name.lower().endswith(".pptx"):
                                            try:
                                                from pptx import Presentation
                                                temp_pptx = f"temp_{file_uuid}.pptx"
                                                blob_client_temp = container_client.get_blob_client(final_blob_name)
                                                with open(temp_pptx, "wb") as f:
                                                    f.write(blob_client_temp.download_blob().readall())
                                                
                                                prs = Presentation(temp_pptx)
                                                font_name = "Times New Roman"
                                                def change_font(shapes):
                                                    for shape in shapes:
                                                        if shape.has_text_frame:
                                                            for p in shape.text_frame.paragraphs:
                                                                for r in p.runs: r.font.name = font_name
                                                        if shape.has_table:
                                                            for r in shape.table.rows:
                                                                for c in r.cells:
                                                                    if c.text_frame:
                                                                        for p in c.text_frame.paragraphs:
                                                                            for run in p.runs: run.font.name = font_name
                                                        if shape.shape_type == 6: change_font(shape.shapes)
                                                
                                                for slide in prs.slides: change_font(slide.shapes)
                                                prs.save(temp_pptx)
                                                
                                                with open(temp_pptx, "rb") as f:
                                                    blob_client_temp.upload_blob(f, overwrite=True)
                                                os.remove(temp_pptx)
                                                st.toast("PPTX Font Fixed")
                                            except Exception as e:
                                                st.warning(f"PPTX Font Fix Failed: {e}")

                                        # Generate Download Link
                                        sas = generate_sas_url(blob_service_client, CONTAINER_NAME, final_blob_name)
                                        download_links.append(f"[{final_file_name} 다운로드]({sas})")
                                    
                                    for link in download_links:
                                        st.markdown(link, unsafe_allow_html=True)

                            except Exception as e:
                                st.error(f"처리 중 오류: {e}")
                            


elif menu == "파일 보관함":
    st.subheader("📂 클라우드 파일 보관함")
    
    # -----------------------------
    # 1. 파일 직접 업로드 (Save)
    # -----------------------------
    with st.expander("📤 파일 직접 업로드 (번역 없이 저장)", expanded=False):
        upload_archive = st.file_uploader("보관함에 저장할 파일 선택", key="archive_upload")
        if st.button("저장하기", disabled=not upload_archive):
            try:
                blob_service_client = get_blob_service_client()
                container_client = blob_service_client.get_container_client(CONTAINER_NAME)
                
                file_uuid = str(uuid.uuid4())[:8]
                blob_name = f"input/{file_uuid}/{upload_archive.name}"
                blob_client = container_client.get_blob_client(blob_name)
                blob_client.upload_blob(upload_archive, overwrite=True)
                st.success(f"'{upload_archive.name}' 업로드 완료!")
                time.sleep(1)
                st.rerun()
            except Exception as e:
                st.error(f"업로드 실패: {e}")

    st.divider()
    
    if st.button("🔄 목록 새로고침"):
        st.rerun()
        
    try:
        blob_service_client = get_blob_service_client()
        container_client = blob_service_client.get_container_client(CONTAINER_NAME)
        
        # 탭으로 Input/Output 구분
        tab1, tab2 = st.tabs(["원본 문서 (Input)", "번역된 문서 (Output)"])
        
        def render_file_list(prefix, tab_name):
            blobs = list(container_client.list_blobs(name_starts_with=prefix))
            blobs.sort(key=lambda x: x.creation_time, reverse=True)
            
            if not blobs:
                st.info(f"{tab_name}에 파일이 없습니다.")
                return

            for i, blob in enumerate(blobs):
                file_name = blob.name.split("/")[-1]
                creation_time = blob.creation_time.strftime('%Y-%m-%d %H:%M')
                
                with st.container():
                    col1, col2, col3 = st.columns([6, 2, 2])
                    
                    with col1:
                        sas_url = generate_sas_url(blob_service_client, CONTAINER_NAME, blob.name)
                        st.markdown(f"**[{file_name}]({sas_url})**")
                        st.caption(f"📅 {creation_time} | 📦 {blob.size / 1024:.1f} KB")
                    
                    with col2:
                        # 수정 (이름 변경)
                        with st.popover("수정"):
                            new_name = st.text_input("새 파일명", value=file_name, key=f"rename_{prefix}_{i}")
                            if st.button("이름 변경", key=f"btn_rename_{prefix}_{i}"):
                                try:
                                    # 새 경로 생성 (UUID 폴더 구조 유지)
                                    path_parts = blob.name.split("/")
                                    # path_parts = ['input', 'uuid', 'filename']
                                    if len(path_parts) >= 3:
                                        new_blob_name = f"{path_parts[0]}/{path_parts[1]}/{new_name}"
                                    else:
                                        # 폴더 구조가 다를 경우 그냥 같은 폴더에
                                        folder = "/".join(path_parts[:-1])
                                        new_blob_name = f"{folder}/{new_name}"
                                    
                                    # 복사 (Rename은 Copy + Delete)
                                    source_blob = container_client.get_blob_client(blob.name)
                                    dest_blob = container_client.get_blob_client(new_blob_name)
                                    
                                    # SAS URL for Copy Source
                                    source_sas = generate_sas_url(blob_service_client, CONTAINER_NAME, blob.name)
                                    
                                    dest_blob.start_copy_from_url(source_sas)
                                    
                                    # 복사 완료 대기 (간단한 폴링)
                                    for _ in range(10):
                                        props = dest_blob.get_blob_properties()
                                        if props.copy.status == "success":
                                            break
                                        time.sleep(0.5)
                                    
                                    # 원본 삭제
                                    source_blob.delete_blob()
                                    st.success("이름 변경 완료!")
                                    time.sleep(1)
                                    st.rerun()
                                except Exception as e:
                                    st.error(f"이름 변경 실패: {e}")

                    with col3:
                        # 삭제
                        if st.button("삭제", key=f"del_{prefix}_{i}", type="secondary"):
                            try:
                                container_client.delete_blob(blob.name)
                                st.success("삭제되었습니다.")
                                time.sleep(1)
                                st.rerun()
                            except Exception as e:
                                st.error(f"삭제 실패: {e}")
                    
                    st.divider()

        with tab1:
            render_file_list("input/", "원본 문서")
            
        with tab2:
            render_file_list("output/", "번역된 문서")
                
    except Exception as e:
        st.error(f"파일 목록을 불러오는 중 오류 발생: {e}")
